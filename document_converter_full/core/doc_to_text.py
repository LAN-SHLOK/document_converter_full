"""
doc_to_text.py

Convert many document formats (image, scanned PDF, text-PDF, DOCX, PPTX) to a .txt file.
Uses PyMuPDF/pdf2image for PDFs, python-docx for DOCX, python-pptx for PPTX,
and two OCR engines (pytesseract & easyocr) with word-level merging by confidence.

Author: ChatGPT (GPT-5 Thinking mini) — revised
"""
# near top of core/doc_to_text.py — add these imports if not present
import shutil
# existing imports: Path, PIL, pytesseract, easyocr, etc.

# determine whether Tesseract binary is present in PATH
TESSERACT_AVAILABLE = shutil.which("tesseract") is not None

# EasyOCR availability (you may already have this variable)
try:
    import easyocr
    HAS_EASYOCR = True
except Exception:
    easyocr = None
    HAS_EASYOCR = False

# create reader lazily
_reader_easyocr = None
def get_easyocr_reader():
    global _reader_easyocr
    if not HAS_EASYOCR:
        return None
    if _reader_easyocr is None:
        _reader_easyocr = easyocr.Reader(['en'], gpu=False)
    return _reader_easyocr
import os
import sys
import re
import tempfile
from pathlib import Path
from typing import List, Tuple, Dict, Optional

import fitz  # PyMuPDF
from pdf2image import convert_from_path
from PIL import Image
import pytesseract
import easyocr
import cv2
import numpy as np
import pandas as pd  # required for pytesseract DATAFRAME output
from docx import Document as DocxDocument
from pptx import Presentation
from spellchecker import SpellChecker

# -------------- Config --------------
# If tesseract binary is not in PATH, set it here (example windows)
# pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

EASYOCR_LANGS = ['en']  # adjust for other languages if needed, e.g. ['en','hi']
DO_SPELL_CORRECTION = False  # set True to enable a light spell-correction pass
# -------------------------------------

# initialize heavy objects once
reader_easyocr = easyocr.Reader(EASYOCR_LANGS, gpu=False)
GLOBAL_SPELLER = SpellChecker() if DO_SPELL_CORRECTION else None


# ---------- Utilities ----------
def preprocess_image_for_ocr(pil_image: Image.Image) -> Image.Image:
    """Preprocess PIL image: convert to grayscale, denoise, adaptive threshold."""
    img = np.array(pil_image.convert('RGB'))
    gray = cv2.cvtColor(img, cv2.COLOR_RGB2GRAY)
    denoised = cv2.fastNlMeansDenoising(gray, h=10)
    th = cv2.adaptiveThreshold(denoised, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                               cv2.THRESH_BINARY, 15, 10)
    return Image.fromarray(th)


def ocr_pytesseract_with_conf(img: Image.Image) -> List[Tuple[str, float]]:
    """
    Run pytesseract and return list of (word, confidence) in reading order.
    Uses TSV / dataframe output to get confidences.
    """
    try:
        tsv = pytesseract.image_to_data(img, output_type=pytesseract.Output.DATAFRAME)
    except Exception:
        # fallback to simple string if data frame unavailable
        raw = pytesseract.image_to_string(img)
        words = [w for w in re.split(r'\s+', raw.strip()) if w]
        return [(w, 0.0) for w in words]

    words_with_conf = []
    if tsv is None or tsv.empty:
        return words_with_conf

    for _, row in tsv.iterrows():
        txt = str(row.get('text', '')).strip()
        # row['conf'] may be numeric, string, or NaN
        conf_val = row.get('conf', None)
        conf = -1.0
        try:
            # use pandas-safe checks
            if pd.isna(conf_val):
                conf = -1.0
            else:
                conf = float(conf_val)
        except Exception:
            try:
                conf = float(str(conf_val))
            except Exception:
                conf = -1.0

        if txt != "" and conf != -1.0:
            words_with_conf.append((txt, max(0.0, min(1.0, conf / 100.0))))  # normalize 0..1
    return words_with_conf


def ocr_easyocr_with_conf(img: Image.Image) -> List[Tuple[str, float]]:
    """
    Run EasyOCR and return list of (word, confidence).
    EasyOCR returns boxes with text; split by whitespace for per-word approach.
    """
    arr = np.array(img.convert('RGB'))
    results = reader_easyocr.readtext(arr)
    words_conf = []
    for bbox, text, conf in results:
        parts = re.split(r'\s+', text.strip())
        for p in parts:
            if p:
                # easyocr confidences are commonly 0..1
                try:
                    conff = float(conf)
                    if conff > 1.0:  # sometimes returned as percent
                        conff = conff / 100.0
                except Exception:
                    conff = 0.0
                words_conf.append((p, max(0.0, min(1.0, conff))))
    return words_conf


def words_to_text(words: List[str]) -> str:
    """Join list of words into a reasonably spaced text string."""
    return " ".join(words)


def merge_word_lists_by_conf(wlist1: List[Tuple[str, float]], wlist2: List[Tuple[str, float]]) -> Tuple[List[str], float]:
    """
    Merge two sequences of (word, conf) by position: for each position take the word with higher confidence.
    If lengths differ, append remaining words from the longer list.
    Returns merged_words_list and average_confidence (0..1).
    """
    merged = []
    confs = []
    n = max(len(wlist1), len(wlist2))
    for i in range(n):
        w1, c1 = (("", 0.0) if i >= len(wlist1) else wlist1[i])
        w2, c2 = (("", 0.0) if i >= len(wlist2) else wlist2[i])
        if c1 >= c2:
            chosen, chosen_conf = w1, c1
        else:
            chosen, chosen_conf = w2, c2
        if chosen == "":
            if w1:
                chosen, chosen_conf = w1, c1
            elif w2:
                chosen, chosen_conf = w2, c2
        if chosen:
            merged.append(chosen)
            confs.append(float(chosen_conf))
    avg_conf = float(np.mean(confs)) if confs else 0.0
    return merged, avg_conf


# ---------- Format specific extractors ----------
def extract_text_from_docx(path: Path) -> str:
    # note: python-docx handles .docx only; .doc (older binary) is not supported
    doc = DocxDocument(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip() != ""]
    return "\n".join(paragraphs)


def extract_text_from_pptx(path: Path) -> str:
    # python-pptx supports .pptx only; .ppt (binary) needs conversion
    prs = Presentation(str(path))
    out_lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = shape.text.strip()
                if txt:
                    out_lines.append(txt)
    return "\n".join(out_lines)


def extract_text_from_text_pdf(path: Path) -> str:
    """
    Attempt text extraction using PyMuPDF (fast and reliable for PDFs with embedded text).
    """
    out = []
    doc = fitz.open(str(path))
    for page in doc:
        ptxt = page.get_text("text")
        if ptxt and ptxt.strip():
            out.append(ptxt.strip())
    return "\n\n".join(out)


def pdf_to_images(path: Path, dpi: int = 300) -> List[Image.Image]:
    # convert_to images via pdf2image (uses poppler)
    pil_pages = convert_from_path(str(path), dpi=dpi)
    return pil_pages


def extract_text_by_ocr_from_image(img: Image.Image) -> Tuple[str, float]:
    """
    Run OCR on a PIL image using:
      1) pytesseract if tesseract binary available
      2) easyocr if tesseract missing but EasyOCR installed
      3) else raise RuntimeError with a clear message for the caller
    Returns (text, avg_confidence)
    """
    # Preprocess
    prep = preprocess_image_for_ocr(img)

    # If tesseract binary exists -> use both pytesseract + easyocr (if available) and merge
    if TESSERACT_AVAILABLE:
        try:
            t_words = ocr_pytesseract_with_conf(prep)
        except Exception:
            t_words = []
        try:
            e_words = ocr_easyocr_with_conf(prep)
        except Exception:
            e_words = []
        # same merging logic as before
        if not t_words and not e_words:
            raw = pytesseract.image_to_string(prep)
            words = [w for w in re.split(r'\s+', raw.strip()) if w]
            return words_to_text(words), 0.0
        if not t_words:
            merged_words = [w for w, _ in e_words]
            avg_conf = float(np.mean([c for _, c in e_words])) if e_words else 0.0
            return words_to_text(merged_words), avg_conf
        if not e_words:
            merged_words = [w for w, _ in t_words]
            avg_conf = float(np.mean([c for _, c in t_words])) if t_words else 0.0
            return words_to_text(merged_words), avg_conf
        merged_words, avg_conf = merge_word_lists_by_conf(t_words, e_words)
        return words_to_text(merged_words), avg_conf

    # If no tesseract binary, but EasyOCR installed, use EasyOCR alone
    if HAS_EASYOCR:
        reader = get_easyocr_reader()
        if reader is None:
            # unexpected: easyocr import succeeded earlier but creating reader failed
            raise RuntimeError("EasyOCR could not be initialized on this host.")
        try:
            e_words = ocr_easyocr_with_conf(prep)
            merged_words = [w for w, _ in e_words]
            avg_conf = float(np.mean([c for _, c in e_words])) if e_words else 0.0
            return words_to_text(merged_words), avg_conf
        except Exception as ex:
            raise RuntimeError(f"EasyOCR failed: {ex}")

    # Neither tesseract nor EasyOCR available — tell caller to use Docker or install Tesseract
    raise RuntimeError(
        "Tesseract binary not found in PATH and EasyOCR is not installed on this host. "
        "Image / scanned-PDF OCR is unavailable. To enable full OCR, either install Tesseract on the server "
        "or deploy using the provided Dockerfile (which installs tesseract & poppler). See README for details."
    )


def extract_text_from_pdf_with_mixed_strategy(path: Path) -> Tuple[str, float]:
    """
    Try native text extraction first. If not enough text (threshold), fall back to OCR on pages.
    Returns combined text and an estimated avg confidence.
    """
    text = extract_text_from_text_pdf(path)
    if len(text.strip()) > 50:
        return text, 0.99

    images = pdf_to_images(path, dpi=300)
    all_texts = []
    confs = []
    for img in images:
        txt, conf = extract_text_by_ocr_from_image(img)
        if txt.strip():
            all_texts.append(txt)
            confs.append(conf)
    combined = "\n\n".join(all_texts)
    avg_conf = float(np.mean(confs)) if confs else 0.0
    return combined, avg_conf


def extract_text_from_image_file(path: Path) -> Tuple[str, float]:
    img = Image.open(str(path))
    return extract_text_by_ocr_from_image(img)


# ---------- Spell correction (optional, light) ----------
def lightly_spell_correct_text(text: str) -> str:
    if GLOBAL_SPELLER is None:
        return text
    # only correct probable words; keep punctuation spacing
    tokens = re.split(r'(\W+)', text)
    corrected_tokens = []
    for tok in tokens:
        if re.match(r'^[A-Za-z]{3,}$', tok):
            corr = GLOBAL_SPELLER.correction(tok)
            corrected_tokens.append(corr if corr else tok)
        else:
            corrected_tokens.append(tok)
    return "".join(corrected_tokens)


# ---------- Important details extractor (simple heuristics) ----------
def extract_important_details(text: str) -> Dict[str, List[str]]:
    details = {"emails": [], "phones": [], "key_values": [], "dates": []}
    # emails
    details['emails'] = list(set(re.findall(r'[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}', text)))
    # phones (loose)
    details['phones'] = list(set(re.findall(r'\+?\d[\d\-\s]{6,}\d', text)))
    # key: value lines (simple)
    kvs = []
    for line in text.splitlines():
        if ':' in line:
            left, right = line.split(':', 1)
            left = left.strip()
            right = right.strip()
            if left and right and len(left) < 80:
                kvs.append({left: right})
    details['key_values'] = kvs
    # date-like tokens (very loose)
    dates = re.findall(r'\b(?:\d{1,2}[/-]\d{1,2}[/-]\d{2,4}|\d{4}-\d{2}-\d{2}|(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{1,2},?\s*\d{0,4})\b', text, flags=re.IGNORECASE)
    details['dates'] = list(set(dates))
    return details


# ---------- Main orchestration ----------
def convert_any_to_text(filepath: str, do_spell_correct: bool = DO_SPELL_CORRECTION) -> Tuple[str, float, str]:
    p = Path(filepath)
    if not p.exists():
        raise FileNotFoundError(filepath)

    suffix = p.suffix.lower()
    combined_text = ""
    estimated_conf = 0.0

    if suffix in ['.docx']:
        try:
            combined_text = extract_text_from_docx(p)
            estimated_conf = 0.999
        except Exception as e:
            print("DOCX extraction failed, falling back to OCR:", e)
            combined_text, estimated_conf = extract_text_from_image_file(p)

    elif suffix in ['.doc']:
        # .doc (binary) not supported by python-docx; recommend conversion externally
        raise RuntimeError(".doc (binary) not supported. Convert to .docx (LibreOffice/pandoc) or use antiword.")

    elif suffix in ['.pptx']:
        try:
            combined_text = extract_text_from_pptx(p)
            estimated_conf = 0.995
        except Exception as e:
            print("PPTX extraction failed, falling back to OCR:", e)
            combined_text, estimated_conf = extract_text_from_image_file(p)

    elif suffix in ['.ppt']:
        raise RuntimeError(".ppt (binary) not supported. Convert to .pptx first (LibreOffice/pandoc).")

    elif suffix in ['.pdf']:
        combined_text, estimated_conf = extract_text_from_pdf_with_mixed_strategy(p)

    elif suffix in ['.jpg', '.jpeg', '.png', '.tiff', '.bmp', '.webp']:
        combined_text, estimated_conf = extract_text_from_image_file(p)

    else:
        try:
            with open(p, 'r', encoding='utf-8') as f:
                combined_text = f.read()
            estimated_conf = 0.999
        except Exception:
            try:
                combined_text, estimated_conf = extract_text_from_image_file(p)
            except Exception as e:
                raise RuntimeError("Unsupported file type and fallback OCR failed: " + str(e))

    if do_spell_correct and combined_text.strip():
        combined_text = lightly_spell_correct_text(combined_text)

    out_path = p.with_suffix('.txt')
    with open(out_path, 'w', encoding='utf-8') as fout:
        fout.write(combined_text)

    return str(out_path), float(estimated_conf), combined_text


# ---------- Command-line entry ----------
def main():
    if len(sys.argv) >= 2:
        infile = sys.argv[1]
    else:
        infile = input("Enter path to document (PDF, image, DOCX, PPTX, etc): ").strip()

    try:
        out_file, conf, _ = convert_any_to_text(infile)
        print(f"Saved text to: {out_file}")
        print(f"Estimated average confidence: {conf*100:.2f}%")
        if conf < 0.5:
            print("Warning: low OCR confidence. Consider providing a cleaner scan, higher DPI, or enabling spell-correction.")
    except Exception as e:
        print("Error:", e)


if __name__ == "__main__":
    main()
